"""
Document parsing + chunking + Qdrant ingestion for the regulatory/policy
corpus (SR 26-2, superseded SR 11-7, internal policy docs, model design
docs). Fully local — pypdf/pdfplumber for parsing, no external OCR service,
per the stack constraint.

Chunking strategy: page-aware, ~800 tokens with 100-token overlap. Page
boundaries are preserved in metadata specifically so a citation can point
to "SR 26-2, page 14" rather than an opaque chunk index — that's what makes
a Finding's evidence_ref legible to a human reviewer, not just resolvable
by code.
"""
import hashlib
import uuid
from dataclasses import dataclass
from pathlib import Path

import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
from qdrant_client import QdrantClient
from qdrant_client.http.models import Distance, PointStruct, VectorParams

from app.core.config import get_settings

settings = get_settings()

CHUNK_SIZE_CHARS = 3200  # roughly 800 tokens
CHUNK_OVERLAP_CHARS = 400


@dataclass
class ParsedPage:
    page_number: int
    text: str


@dataclass
class Chunk:
    chunk_id: str
    text: str
    page_number: int
    source_document: str
    doc_hash: str


def parse_pdf(path: Path) -> list[ParsedPage]:
    pages: list[ParsedPage] = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            pages.append(ParsedPage(page_number=i, text=text))
    return pages


def parse_docx(path: Path) -> list[ParsedPage]:
    doc = DocxDocument(str(path))
    full_text = "\n".join(p.text for p in doc.paragraphs)
    return [ParsedPage(page_number=1, text=full_text)]


def parse_pptx(path: Path) -> list[ParsedPage]:
    prs = Presentation(str(path))
    pages: list[ParsedPage] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = [shape.text for shape in slide.shapes if shape.has_text_frame]
        pages.append(ParsedPage(page_number=i, text="\n".join(texts)))
    return pages


def parse_document(path: Path) -> list[ParsedPage]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return parse_pdf(path)
    if suffix == ".docx":
        return parse_docx(path)
    if suffix == ".pptx":
        return parse_pptx(path)
    if suffix in (".txt", ".md"):
        return [ParsedPage(page_number=1, text=path.read_text(encoding="utf-8"))]
    raise ValueError(f"Unsupported document type: {suffix}")


def chunk_pages(pages: list[ParsedPage], source_document: str) -> list[Chunk]:
    doc_hash = hashlib.sha256(source_document.encode()).hexdigest()[:16]
    chunks: list[Chunk] = []
    for page in pages:
        text = page.text
        if not text.strip():
            continue
        start = 0
        while start < len(text):
            end = min(start + CHUNK_SIZE_CHARS, len(text))
            chunk_text = text[start:end]
            chunk_id = f"{doc_hash}-p{page.page_number}-{start}"
            chunks.append(
                Chunk(
                    chunk_id=chunk_id,
                    text=chunk_text,
                    page_number=page.page_number,
                    source_document=source_document,
                    doc_hash=doc_hash,
                )
            )
            if end == len(text):
                break
            start = end - CHUNK_OVERLAP_CHARS
    return chunks


class RegulatoryCorpusIndex:
    """Thin wrapper around Qdrant for the regulatory/policy collection.
    Embedding model choice is deferred to Day 2 (wired alongside the
    gateway); this class assumes an embed_fn is passed in so ingestion.py
    has no hidden dependency on which embedding provider is active."""

    def __init__(self, client: QdrantClient | None = None):
        self.client = client or QdrantClient(url=settings.qdrant_url)
        self.collection = settings.qdrant_collection_regulatory

    def ensure_collection(self, vector_size: int = 1024) -> None:
        existing = [c.name for c in self.client.get_collections().collections]
        if self.collection not in existing:
            self.client.create_collection(
                collection_name=self.collection,
                vectors_config=VectorParams(size=vector_size, distance=Distance.COSINE),
            )

    def upsert_chunks(self, chunks: list[Chunk], embed_fn) -> int:
        if not chunks:
            return 0
        vectors = embed_fn([c.text for c in chunks])
        points = [
            PointStruct(
                id=str(uuid.uuid5(uuid.NAMESPACE_URL, chunk.chunk_id)),
                vector=vector,
                payload={
                    "text": chunk.text,
                    "page_number": chunk.page_number,
                    "source_document": chunk.source_document,
                    "doc_hash": chunk.doc_hash,
                },
            )
            for chunk, vector in zip(chunks, vectors)
        ]
        self.client.upsert(collection_name=self.collection, points=points)
        return len(points)
